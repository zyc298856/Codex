from pathlib import Path
import sys
from pptx import Presentation


def main() -> int:
    if len(sys.argv) != 2:
        print("usage: inspect_pptx.py <deck.pptx>")
        return 2
    ppt = Path(sys.argv[1])
    prs = Presentation(str(ppt))
    print(f"exists={ppt.exists()} size_kb={ppt.stat().st_size / 1024:.1f} slides={len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        pics = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text_frame.text.strip()
                if value:
                    texts.append(value.replace("\n", " / ")[:90])
            if shape.shape_type == 13:
                pics += 1
        print(f"[{i}] text_items={len(texts)} pictures={pics} first={texts[:2]}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
