from __future__ import annotations

from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[2]
VENDOR = ROOT.parent.parent / ".vendor"
if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "paper"
OUT_PPTX = OUT_DIR / "defense_presentation_v1.pptx"
OUT_NOTES = OUT_DIR / "defense_presentation_speaker_notes_v1.md"

W, H = Inches(13.333), Inches(7.5)

COLORS = {
    "navy": RGBColor(14, 31, 50),
    "blue": RGBColor(18, 102, 170),
    "cyan": RGBColor(0, 153, 204),
    "teal": RGBColor(38, 166, 154),
    "green": RGBColor(74, 173, 92),
    "orange": RGBColor(239, 137, 43),
    "red": RGBColor(214, 71, 63),
    "ink": RGBColor(30, 41, 59),
    "muted": RGBColor(100, 116, 139),
    "light": RGBColor(241, 247, 252),
    "panel": RGBColor(229, 239, 248),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(196, 212, 226),
}


def set_run(run, size=24, color="ink", bold=False, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = COLORS[color]
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=24, color="ink", bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Microsoft YaHei",
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, font=font)
    return box


def add_multiline(slide, lines, x, y, w, h, size=22, color="ink", bullet=False, gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.text = "• " + line
            for run in p.runs:
                set_run(run, size=size, color=color)
        else:
            run = p.add_run()
            run.text = line
            set_run(run, size=size, color=color)
    return box


def add_title(slide, title, subtitle=None, section=None, dark=False):
    color = "white" if dark else "navy"
    if section:
        add_text(slide, section, 0.64, 0.28, 4.0, 0.32, 12, "cyan" if not dark else "teal", True)
    add_text(slide, title, 0.62, 0.55, 8.6, 0.62, 28, color, True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.13, 10.8, 0.36, 14, "muted" if not dark else "panel")


def add_footer(slide, page):
    add_text(slide, "基于嵌入式平台的目标检测系统研究", 0.62, 7.04, 5.4, 0.22, 8.5, "muted")
    add_text(slide, f"{page:02d}", 12.08, 7.03, 0.48, 0.22, 9, "muted", align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(6.91), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def rect(slide, x, y, w, h, fill="panel", line="line", radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = COLORS[fill]
    s.line.color.rgb = COLORS[line]
    s.line.width = Pt(0.8)
    return s


def chip(slide, text, x, y, w, fill="blue"):
    s = rect(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    add_text(slide, text, x + 0.08, y + 0.06, w - 0.16, 0.18, 9.5, "white", True, align=PP_ALIGN.CENTER)
    return s


def arrow(slide, x1, y1, x2, y2, color="blue", width=1.8):
    a = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    a.line.color.rgb = COLORS[color]
    a.line.width = Pt(width)
    try:
        a.line.end_arrowhead = True
    except Exception:
        pass
    return a


def node(slide, title, x, y, w, h, fill="white", accent="blue", small=None):
    rect(slide, x, y, w, h, fill=fill, line="line", radius=True)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS[accent]
    accent_bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.32, 0.28, 14, "navy", True)
    if small:
        add_text(slide, small, x + 0.18, y + 0.5, w - 0.32, h - 0.56, 10.5, "muted")


def metric(slide, label, value, note, x, y, w, color="blue"):
    add_text(slide, value, x, y, w, 0.48, 30, color, True)
    add_text(slide, label, x, y + 0.54, w, 0.24, 12, "navy", True)
    add_text(slide, note, x, y + 0.84, w, 0.42, 9.5, "muted")


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    slides = []

    # 1 Cover
    s = prs.slides.add_slide(blank)
    slides.append(s)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    rect(s, 0, 0, 13.333, 7.5, fill="navy", line="navy")
    rect(s, 8.2, -0.1, 5.2, 7.8, fill="blue", line="blue")
    rect(s, 8.85, 0.48, 3.7, 5.7, fill="light", line="light")
    add_text(s, "RK3588", 9.08, 0.9, 2.4, 0.5, 30, "blue", True, font="Aptos Display")
    add_text(s, "NPU 多线程\n实时目标检测", 9.08, 1.56, 3.0, 1.05, 25, "navy", True)
    for i, label in enumerate(["PyTorch", "ONNX", "RKNN", "RTSP"]):
        chip(s, label, 9.12, 3.0 + i * 0.52, 1.2, fill="teal" if i != 2 else "orange")
    add_text(s, "基于嵌入式平台的目标检测系统研究", 0.74, 1.2, 6.9, 0.78, 31, "white", True)
    add_text(s, "本科毕业论文（设计）答辩汇报", 0.78, 2.1, 4.8, 0.32, 17, "panel")
    add_text(s, "朱奕澄  |  通信与信息工程学院  |  电子信息工程", 0.8, 5.85, 6.6, 0.28, 13, "panel")
    add_text(s, "研究重点：模型迁移、RKNN 部署、实时流水线、NPU 多 context 并行与硬件优化实验", 0.8, 6.28, 7.0, 0.3, 12, "teal")

    # 2 Background
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "研究背景：边缘端小目标检测需要系统级优化", "无人机/小目标检测不仅是模型精度问题，也是嵌入式平台上的实时工程问题", "01 背景")
    add_multiline(s, [
        "远距离目标在图像中占比小，容易受天空、建筑、树木、鸟类和光照干扰。",
        "嵌入式平台算力、内存带宽、视频编码和线程调度共同影响实时显示效果。",
        "导师关注的核心方向是 RK3588 NPU 多线程与硬件协同优化，而不是单纯换一个更大的模型。"
    ], 0.75, 1.75, 6.1, 2.1, 18, "ink", bullet=True, gap=5)
    node(s, "检测难点", 7.45, 1.52, 4.9, 0.92, small="小目标、背景干扰、框抖动、漏检/误检")
    node(s, "部署难点", 7.45, 2.68, 4.9, 0.92, accent="teal", small="NPU 推理、CPU 后处理、内存拷贝、RTSP 输出")
    node(s, "论文定位", 7.45, 3.84, 4.9, 0.92, accent="orange", small="以工程实现和实验验证为主，稳健描述边界")
    add_footer(s, 2)

    # 3 Objectives and completion
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "任务目标与完成情况", "以“能跑通、可复现、可解释”为主线推进", "02 目标")
    headers = ["任务线", "当前状态", "答辩表述"]
    xs = [0.8, 3.25, 6.55]
    ws = [2.0, 2.55, 5.75]
    y0 = 1.55
    for x, w, h in zip(xs, ws, headers):
        rect(s, x, y0, w, 0.46, fill="blue", line="blue")
        add_text(s, h, x + 0.08, y0 + 0.11, w - 0.16, 0.16, 10.5, "white", True, align=PP_ALIGN.CENTER)
    rows = [
        ("模型迁移", "已完成", "PyTorch→ONNX→RKNN，稳定主线采用 FP RKNN"),
        ("实时检测", "已跑通", "固定视频和 USB 摄像头 RTSP 均可输出检测画面"),
        ("NPU 多线程", "已验证", "双 context 可作为硬件并行实验重点"),
        ("RGA/Zero-copy", "实验路径", "已有可切换实现，作为硬件优化对比而非默认路径"),
        ("GPIO 报警", "替代实现", "无外设条件下使用软件告警叠加和日志记录"),
    ]
    for i, row in enumerate(rows):
        y = y0 + 0.5 + i * 0.72
        fill = "light" if i % 2 == 0 else "white"
        for x, w, val in zip(xs, ws, row):
            rect(s, x, y, w, 0.64, fill=fill, line="line")
            add_text(s, val, x + 0.1, y + 0.14, w - 0.2, 0.18, 10.8, "ink", bold=(val in ["已完成", "已跑通", "已验证"]))
    add_footer(s, 3)

    # 4 Architecture
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "系统总体架构", "从视频输入到检测结果输出采用分层、可切换的流水线设计", "03 架构")
    node(s, "视频输入", 0.85, 2.0, 1.55, 0.78, accent="teal", small="文件 / USB / RTSP")
    node(s, "预处理", 2.9, 2.0, 1.55, 0.78, accent="cyan", small="OpenCV / RGA 可选")
    node(s, "RKNN 推理", 4.95, 2.0, 1.65, 0.78, accent="orange", small="NPU context")
    node(s, "后处理", 7.1, 2.0, 1.55, 0.78, accent="blue", small="解码、阈值、NMS")
    node(s, "策略控制", 9.15, 2.0, 1.65, 0.78, accent="green", small="间隔检测、ROI、跟踪")
    node(s, "输出", 11.3, 2.0, 1.45, 0.78, accent="red", small="视频 / CSV / RTSP")
    for x1, x2 in [(2.4, 2.88), (4.45, 4.93), (6.6, 7.08), (8.65, 9.13), (10.8, 11.28)]:
        arrow(s, x1, 2.39, x2, 2.39, "blue")
    add_multiline(s, [
        "固定视频路径：用于可重复实验、生成带框视频、CSV、ROI JSONL 和告警事件。",
        "实时 RTSP 路径：用于摄像头演示和实时观看，强调低延迟与视觉稳定性。",
        "硬件优化路径：multi-context、RGA、zero-copy、profiling 均通过环境变量开关控制。"
    ], 1.05, 4.0, 11.2, 1.35, 16, "ink", bullet=True, gap=4)
    add_footer(s, 4)

    # 5 Model migration
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "模型迁移：从训练模型到 RK3588 可执行模型", "核心不是只导出模型，而是保证输出格式和后处理逻辑一致", "04 模型")
    for i, (t, st) in enumerate([
        ("训练权重", "best.pt"),
        ("中间格式", "best.onnx"),
        ("板端模型", "best.rk3588.fp.rknn"),
        ("C++ 后处理", "[1,5,8400] / end2end 兼容")
    ]):
        x = 0.9 + i * 3.0
        node(s, t, x, 1.85, 2.25, 0.92, accent=["blue", "cyan", "orange", "green"][i], small=st)
        if i < 3:
            arrow(s, x + 2.25, 2.31, x + 2.72, 2.31)
    rect(s, 1.1, 3.55, 11.05, 1.35, fill="light", line="line", radius=True)
    add_text(s, "输出张量设计", 1.35, 3.82, 1.9, 0.24, 15, "navy", True)
    add_text(s, "8400 = 80×80 + 40×40 + 20×20，来自 YOLOv10 三个检测头；当前单类模型采用 5 通道输出，后处理同时兼容 COCO 84 通道与 end2end 6 通道格式。", 3.0, 3.76, 8.65, 0.52, 15, "ink")
    add_multiline(s, ["稳定主线采用 FP RKNN；INT8 已作为后续量化闭环方向，而不是当前主结论。"], 1.2, 5.45, 10.4, 0.5, 16, "muted", bullet=True)
    add_footer(s, 5)

    # 6 RK3588 platform
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "部署平台：RK3588 的异构计算资源", "CPU、NPU、RGA、视频编解码器共同决定系统实时性", "05 平台")
    metric(s, "NPU 算力", "6 TOPS", "INT8 峰值能力，当前稳定主线为 FP RKNN", 0.95, 1.72, 2.25, "orange")
    metric(s, "CPU 架构", "4+4", "Cortex-A76 + Cortex-A55 big.LITTLE", 3.55, 1.72, 2.45, "blue")
    metric(s, "图形/视频", "Mali-G610", "配合硬件编解码和 RGA 减轻 CPU 负担", 6.25, 1.72, 2.75, "teal")
    metric(s, "系统目标", "实时", "在低功耗板端完成检测和显示输出", 9.45, 1.72, 2.35, "green")
    add_multiline(s, [
        "NPU 负责神经网络推理；CPU 负责线程调度、后处理、日志和策略控制。",
        "RGA 可承担缩放、色彩转换等图像处理任务，但是否提升端到端性能需要实测。",
        "内存带宽与数据拷贝路径会影响推理前后的整体延迟。"
    ], 0.98, 4.05, 11.2, 1.3, 16, "ink", bullet=True, gap=4)
    add_footer(s, 6)

    # 7 Runtime paths
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "两条运行路径：固定视频实验与实时 RTSP 演示", "固定视频保证可复现，RTSP 证明真实部署能力", "06 运行")
    node(s, "rk_yolo_video", 1.1, 1.75, 4.8, 1.0, accent="blue", small="输入视频文件，输出带框视频、CSV、ROI JSONL、告警 CSV")
    node(s, "rk_yolo_live_rtsp", 7.15, 1.75, 4.8, 1.0, accent="teal", small="输入 USB 摄像头或文件，输出实时 RTSP 画面")
    arrow(s, 5.92, 2.25, 7.1, 2.25, "orange", 2.0)
    add_multiline(s, [
        "论文主实验：优先使用固定视频路径，便于重复比较不同参数和硬件优化策略。",
        "答辩演示：优先使用实时 RTSP 路径，可在电脑端观看带框检测画面。",
        "两条路径共享 RKNN 检测器实现，减少代码分叉和实验口径不一致。"
    ], 1.15, 4.0, 11.0, 1.45, 16, "ink", bullet=True, gap=4)
    add_footer(s, 7)

    # 8 Pipeline
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "实时视频处理流水线", "采集、推理、策略控制和输出解耦，避免单个慢模块拖垮系统", "07 流水线")
    y = 1.85
    for i, (t, st, acc) in enumerate([
        ("采集线程", "读取摄像头/视频帧", "teal"),
        ("任务队列", "有限队列，丢弃过期帧", "cyan"),
        ("推理 worker", "独立 RKNN context", "orange"),
        ("结果队列", "按时间戳选择最新结果", "blue"),
        ("显示/推流", "画框、告警、RTSP", "green"),
    ]):
        x = 0.75 + i * 2.45
        node(s, t, x, y, 1.85, 0.95, accent=acc, small=st)
        if i < 4:
            arrow(s, x + 1.85, y + 0.47, x + 2.33, y + 0.47)
    rect(s, 1.15, 4.15, 10.9, 1.0, fill="light", line="line", radius=True)
    add_text(s, "关键设计", 1.38, 4.4, 1.2, 0.25, 15, "navy", True)
    add_text(s, "用“最新帧优先 + 有限队列 + 结果复用”控制实时延迟；固定视频负责严谨对比，实时 RTSP 负责演示。", 2.45, 4.38, 8.75, 0.3, 15, "ink")
    add_footer(s, 8)

    # 9 Multi-context
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "NPU 多 context 并行推理", "这是导师关注的硬件优化主线，也是本项目最强的系统性贡献之一", "08 多线程")
    node(s, "FrameTask 队列", 1.0, 2.0, 2.2, 0.85, accent="blue", small="待检测帧")
    node(s, "Worker 1\nRKNN Context 1", 4.35, 1.38, 2.35, 0.78, accent="orange")
    node(s, "Worker 2\nRKNN Context 2", 4.35, 2.45, 2.35, 0.78, accent="orange")
    node(s, "可扩展 Worker 3/4", 4.35, 3.52, 2.35, 0.78, accent="muted")
    node(s, "DetectionResult", 8.0, 2.45, 2.2, 0.85, accent="green", small="结果合并")
    node(s, "RTSP 输出", 10.95, 2.45, 1.55, 0.85, accent="teal", small="最新结果")
    arrow(s, 3.2, 2.42, 4.32, 1.75)
    arrow(s, 3.2, 2.42, 4.32, 2.84)
    arrow(s, 3.2, 2.42, 4.32, 3.9)
    arrow(s, 6.7, 1.78, 7.98, 2.75)
    arrow(s, 6.7, 2.84, 7.98, 2.84)
    arrow(s, 6.7, 3.9, 7.98, 2.92)
    arrow(s, 10.2, 2.86, 10.93, 2.86)
    add_text(s, "实测结论", 1.05, 5.2, 1.25, 0.25, 15, "navy", True)
    add_text(s, "双 context 在每帧检测模式下提高 NPU FPS 并降低延迟；但实时观看默认仍采用策略优化路径，以获得更稳的画面体验。", 2.15, 5.18, 9.8, 0.3, 15, "ink")
    add_footer(s, 9)

    # 10 Realtime strategy
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "实时策略：检测间隔、ROI、跟踪与平滑", "目标是减少不必要推理，同时让显示框更稳定", "09 策略")
    cols = [
        ("detect_every_n", "每 N 帧执行一次完整 NPU 检测；中间帧复用或跟踪结果。"),
        ("动态 ROI", "围绕历史目标区域裁剪并周期性全帧刷新，兼顾速度和漏检风险。"),
        ("轻量跟踪", "motion 线性位移预测；optflow 使用 Lucas-Kanade 稀疏光流。"),
        ("框平滑", "对检测框坐标做加权融合，减少左右抖动。"),
    ]
    for i, (t, body) in enumerate(cols):
        x = 0.8 + (i % 2) * 6.0
        y = 1.55 + (i // 2) * 1.75
        node(s, t, x, y, 4.95, 1.1, accent=["blue", "teal", "orange", "green"][i], small=body)
    rect(s, 1.1, 5.45, 10.9, 0.6, fill="panel", line="line", radius=True)
    add_text(s, "推荐演示：detect_every_n=2，motion tracking，box smoothing on，dynamic ROI on。", 1.35, 5.62, 10.3, 0.22, 14, "navy", True)
    add_footer(s, 10)

    # 11 Results visual
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "检测效果与可视化输出", "固定视频和实时 RTSP 均支持检测框、类别置信度和告警叠加", "10 效果")
    rect(s, 0.95, 1.55, 5.8, 3.55, fill="navy", line="navy")
    rect(s, 1.28, 1.95, 5.14, 2.72, fill="light", line="light")
    # Mock detection frame drawn with editable shapes.
    for i in range(6):
        rect(s, 1.45 + i * 0.72, 2.18, 0.24, 1.9, fill="panel", line="panel")
    drone = slide_shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.36), Inches(3.08), Inches(0.92), Inches(0.25))
    drone.fill.solid()
    drone.fill.fore_color.rgb = COLORS["ink"]
    drone.line.fill.background()
    bbox = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.22), Inches(2.88), Inches(1.22), Inches(0.7))
    bbox.fill.background()
    bbox.line.color.rgb = COLORS["green"]
    bbox.line.width = Pt(2.2)
    rect(s, 3.22, 2.55, 1.7, 0.35, fill="green", line="green")
    add_text(s, "drone 0.61", 3.31, 2.63, 1.45, 0.12, 11, "white", True)
    rect(s, 1.28, 1.95, 1.75, 0.32, fill="red", line="red")
    add_text(s, "UAV ALERT", 1.45, 2.03, 1.35, 0.12, 10, "white", True)
    add_text(s, "示意：检测框 + 告警条", 2.55, 4.78, 2.6, 0.22, 11, "panel")
    add_multiline(s, [
        "离线视频：输出带框 MP4、逐帧 CSV、ROI JSONL 和 alarm_events.csv。",
        "实时演示：电脑端可通过 RTSP 地址查看检测画面。",
        "无继电器/蜂鸣器条件下，软件告警提供更直观、可记录的替代反馈。"
    ], 7.25, 1.8, 4.9, 2.1, 16, "ink", bullet=True, gap=5)
    add_footer(s, 11)

    # 12 Experiment comparison
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "实验对比：优化不是“硬件越多越好”", "最终推荐基于端到端效果，而不是单个算子是否硬件加速", "11 实验")
    comparison = [
        ("单 context 每帧检测", "NPU FPS≈8.59\n延迟≈268.89 ms", "硬件基线"),
        ("双 context 每帧检测", "NPU FPS≈10.64\n延迟≈173.58 ms", "并行推理证据"),
        ("策略优化实时显示", "detect_every_n=2\n跟踪+平滑+ROI", "演示默认"),
        ("RGA/zero-copy", "部分阶段有效\n需端到端判断", "优化实验"),
    ]
    for i, (name, val, tag) in enumerate(comparison):
        x = 0.8 + i * 3.05
        rect(s, x, 1.7, 2.55, 2.55, fill="light" if i != 1 else "panel", line="line", radius=True)
        chip(s, tag, x + 0.28, 1.95, 1.55, fill=["blue", "orange", "green", "teal"][i])
        add_text(s, name, x + 0.28, 2.48, 2.1, 0.38, 15, "navy", True)
        add_text(s, val, x + 0.28, 3.08, 2.1, 0.68, 17, "ink", True)
    add_text(s, "论文结论口径", 1.0, 5.15, 1.55, 0.25, 15, "navy", True)
    add_text(s, "多 context 证明 NPU 并行能力；实时演示采用更稳的策略优化配置；RGA 与 zero-copy 作为可切换实验路径记录正负结果。", 2.45, 5.13, 9.7, 0.3, 15, "ink")
    add_footer(s, 12)

    # 13 Profiling and hardware paths
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "性能剖析与硬件优化路径", "通过环境变量开关做可重复实验，避免把探索路径写成默认结论", "12 优化")
    node(s, "RK_YOLO_PROFILE", 0.95, 1.55, 3.3, 0.82, accent="blue", small="prepare / input / run / output / decode / render")
    node(s, "RK_YOLO_ZERO_COPY_INPUT", 0.95, 2.65, 3.3, 0.82, accent="orange", small="rknn_create_mem + rknn_set_io_mem")
    node(s, "RK_YOLO_PREPROCESS=rga", 0.95, 3.75, 3.3, 0.82, accent="teal", small="RGA resize / cvt_resize / letterbox")
    rect(s, 5.25, 1.55, 6.95, 3.02, fill="light", line="line", radius=True)
    add_text(s, "当前判断", 5.55, 1.88, 1.2, 0.25, 15, "navy", True)
    add_multiline(s, [
        "zero-copy 输入上传阶段变快，但整体 `rknn_run` 上升，未作为默认。",
        "RGA 在高分辨率输入和部分预处理路径上有收益，但仍需按端到端延迟判断。",
        "INT8 量化需要校准集、精度对比和板端闭环，当前作为后续优化。"
    ], 5.55, 2.35, 6.1, 1.35, 14, "ink", bullet=True, gap=3)
    add_footer(s, 13)

    # 14 Limitations and adjustment
    s = prs.slides.add_slide(blank)
    slides.append(s)
    add_title(s, "目标调整与不足说明", "答辩时建议正面说明：哪些已完成，哪些保留为后续工作", "13 边界")
    rows = [
        ("INT8 量化", "已做离线准备，缺少完整精度-速度闭环", "后续补充校准集和板端对比"),
        ("RGA 主路径", "已有多条实验路径，默认仍选择稳定 OpenCV/FP 路线", "继续优化 MPP/RGA 数据通路"),
        ("GPIO 报警", "无继电器/蜂鸣器硬件", "软件告警叠加 + alarm CSV 作为可视替代"),
        ("真实飞行场景", "受场地与无人机条件限制", "公开视频 + 教室小目标演示补充验证"),
    ]
    for i, (a, b, c) in enumerate(rows):
        y = 1.45 + i * 0.95
        rect(s, 0.85, y, 2.15, 0.72, fill="panel", line="line")
        rect(s, 3.0, y, 4.55, 0.72, fill="white", line="line")
        rect(s, 7.55, y, 4.65, 0.72, fill="light", line="line")
        add_text(s, a, 1.02, y + 0.2, 1.8, 0.16, 12, "navy", True)
        add_text(s, b, 3.18, y + 0.18, 4.15, 0.2, 11.5, "ink")
        add_text(s, c, 7.75, y + 0.18, 4.1, 0.2, 11.5, "ink")
    add_footer(s, 14)

    # 15 Summary
    s = prs.slides.add_slide(blank)
    slides.append(s)
    rect(s, 0, 0, 13.333, 7.5, fill="navy", line="navy")
    add_text(s, "总结", 0.8, 0.85, 1.5, 0.48, 32, "white", True)
    add_text(s, "本课题完成了从模型到 RK3588 板端实时检测系统的工程闭环，并围绕 NPU 多线程与硬件优化做了可复现实验。", 0.82, 1.55, 9.8, 0.62, 21, "panel", True)
    node(s, "系统闭环", 0.95, 3.05, 3.45, 1.0, fill="light", accent="teal", small="训练/转换/部署/固定视频/RTSP 显示")
    node(s, "优化闭环", 4.95, 3.05, 3.45, 1.0, fill="light", accent="orange", small="detect_every_n、tracking、ROI、multi-context、RGA")
    node(s, "论文闭环", 8.95, 3.05, 3.45, 1.0, fill="light", accent="green", small="实验记录、代码映射、格式合规、附录材料")
    add_text(s, "答辩主线：我不是提出一个全新的检测网络，而是在 RK3588 上完成了可运行、可测量、可解释的嵌入式目标检测系统。", 0.95, 5.45, 11.0, 0.36, 16, "teal", True)
    add_text(s, "谢谢各位老师", 0.95, 6.25, 4.0, 0.36, 19, "white", True)

    # Add footer to non-cover/non-final slides done already.
    prs.save(OUT_PPTX)
    return len(slides)


NOTES = """# 答辩 PPT 讲稿提示（v1）

## 1. 封面
一句话开场：本课题不是单纯训练模型，而是完成目标检测算法从 PC 训练到 RK3588 板端实时运行的移植、优化和验证。

## 2. 研究背景
强调小目标和边缘端实时性的双重困难。自然引出导师关注的 NPU 多线程和硬件优化。

## 3. 目标与完成情况
建议主动说明“已完成/实验路径/替代实现”的边界，显得诚实且专业。

## 4. 系统架构
按视频输入、预处理、RKNN 推理、后处理、策略控制、输出六层讲。重点说两条工具链共享检测器。

## 5. 模型迁移
讲清楚 PyTorch→ONNX→RKNN，以及为什么后处理必须兼容不同输出张量。

## 6. RK3588 平台
强调 NPU、CPU、RGA、视频编解码和内存带宽共同影响实时性。

## 7. 两条运行路径
固定视频负责论文实验复现，实时 RTSP 负责答辩演示和部署效果。

## 8. 实时流水线
说明为什么不用单线程串行：会被慢模块阻塞，导致画面延迟和卡顿。

## 9. NPU 多 context
这是答辩重点。说清楚每个 worker 有独立 RKNN context，避免共享上下文带来的同步冲突。

## 10. 实时策略
解释框抖动如何通过检测间隔、运动预测、光流候选、动态 ROI 和平滑缓解。

## 11. 检测效果
如果现场有板子和视频，可以在这一页后切到 VLC 或输出视频演示。

## 12. 实验对比
答辩口径：multi-context 是硬件并行证据；策略优化是实时观看默认方案。

## 13. 性能剖析
强调没有盲目追求硬件加速，而是用 profiling 记录每一阶段耗时。

## 14. 不足与展望
正面说明 INT8、RGA 主路径、GPIO、真实飞行测试的限制和后续方案。

## 15. 总结
用“系统闭环、优化闭环、论文闭环”收束，最后回到课题贡献。
"""


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    count = build_deck()
    OUT_NOTES.write_text(NOTES, encoding="utf-8")
    print(f"created {OUT_PPTX} ({count} slides)")
    print(f"created {OUT_NOTES}")
