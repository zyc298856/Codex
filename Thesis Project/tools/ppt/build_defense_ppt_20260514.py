from __future__ import annotations

from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[2]
VENDOR = ROOT.parent.parent / ".vendor"
if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "paper" / "defense_ppt"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "朱奕澄+答辩PPT.pptx"
OUT_NOTES = OUT_DIR / "朱奕澄+答辩PPT_5分钟讲稿.md"

IMG_DEMO = ROOT / "eval_runs" / "board_fixed_video_demo_20260513_pexels" / "max_score_detection_frame_200.jpg"
IMG_FIRST = ROOT / "eval_runs" / "board_fixed_video_demo_20260513_pexels" / "first_detection_frame_001.jpg"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

COLORS = {
    "ink": RGBColor(22, 31, 44),
    "muted": RGBColor(86, 99, 116),
    "line": RGBColor(187, 202, 216),
    "white": RGBColor(255, 255, 255),
    "navy": RGBColor(15, 38, 62),
    "deep": RGBColor(21, 52, 82),
    "blue": RGBColor(28, 111, 181),
    "cyan": RGBColor(14, 150, 198),
    "green": RGBColor(52, 158, 99),
    "orange": RGBColor(225, 132, 42),
    "red": RGBColor(203, 54, 61),
    "pale": RGBColor(241, 247, 251),
    "pale2": RGBColor(232, 242, 249),
    "cream": RGBColor(252, 247, 239),
}


def rgb(name: str) -> RGBColor:
    return COLORS[name]


def set_run(run, size=16, color="ink", bold=False, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold


def text(slide, value, x, y, w, h, size=16, color="ink", bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    set_run(run, size=size, color=color, bold=bold, font=font)
    return box


def bullets(slide, lines, x, y, w, h, size=15.2, color="ink", gap=4, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        set_run(run, size=size, color=color)
    return box


def shape(slide, x, y, w, h, fill="pale", line="line", rounded=False, width=0.8):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line)
    s.line.width = Pt(width)
    return s


def pill(slide, label, x, y, w, color="blue", size=10.5):
    shape(slide, x, y, w, 0.34, fill=color, line=color, rounded=True)
    text(slide, label, x + 0.03, y + 0.075, w - 0.06, 0.12, size=size, color="white",
         bold=True, align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color="blue", width=1.9):
    a = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    a.line.color.rgb = rgb(color)
    a.line.width = Pt(width)
    return a


def title(slide, section, main, sub=None):
    text(slide, section, 0.62, 0.30, 3.2, 0.25, size=10.5, color="cyan", bold=True)
    text(slide, main, 0.60, 0.58, 10.8, 0.55, size=28, color="navy", bold=True)
    if sub:
        text(slide, sub, 0.62, 1.14, 11.5, 0.30, size=12.5, color="muted")


def footer(slide, n):
    shape(slide, 0.62, 6.90, 12.05, 0.012, fill="line", line="line")
    text(slide, "基于嵌入式平台的目标检测系统研究", 0.62, 7.02, 5.0, 0.18, size=8.4, color="muted")
    text(slide, f"{n:02d}", 12.18, 7.02, 0.45, 0.18, size=8.5, color="muted", align=PP_ALIGN.RIGHT)


def node(slide, head, body, x, y, w, h, accent="blue", fill="white"):
    shape(slide, x, y, w, h, fill=fill, line="line", rounded=True)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    text(slide, head, x + 0.17, y + 0.13, w - 0.25, 0.22, size=12.4, color="navy", bold=True)
    text(slide, body, x + 0.17, y + 0.43, w - 0.25, h - 0.45, size=9.7, color="muted")


def metric(slide, value, label, x, y, w=1.85, color="blue"):
    text(slide, value, x, y, w, 0.42, size=26, color=color, bold=True, align=PP_ALIGN.CENTER, font="Aptos Display")
    text(slide, label, x, y + 0.48, w, 0.28, size=10.6, color="muted", align=PP_ALIGN.CENTER)


def table(slide, headers, rows, x, y, widths, row_h=0.44, font_size=9.2):
    total = sum(widths)
    shape(slide, x, y, total, 0.02, fill="navy", line="navy")
    yy = y + 0.03
    for i, h in enumerate(headers):
        xx = x + sum(widths[:i])
        shape(slide, xx, yy, widths[i], 0.38, fill="pale2", line="pale2")
        text(slide, h, xx + 0.04, yy + 0.09, widths[i] - 0.08, 0.12, size=font_size, color="navy", bold=True, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y + 0.43 + r * row_h
        for i, c in enumerate(row):
            xx = x + sum(widths[:i])
            text(slide, str(c), xx + 0.05, yy + 0.10, widths[i] - 0.10, 0.16, size=font_size, color="ink",
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)
    shape(slide, x, y + 0.43 + len(rows) * row_h + 0.03, total, 0.02, fill="navy", line="navy")


def add_image(slide, path: Path, x, y, w, h):
    if path.exists() and path.stat().st_size > 0:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        shape(slide, x, y, w, h, fill="pale", line="line", rounded=True)
        text(slide, "图片待替换", x, y + h / 2 - 0.12, w, 0.24, size=14, color="muted", align=PP_ALIGN.CENTER)


def slide_cover(prs, blank):
    s = prs.slides.add_slide(blank)
    shape(s, 0, 0, 13.333, 7.5, fill="navy", line="navy")
    shape(s, 8.30, 0, 5.05, 7.5, fill="blue", line="blue")
    shape(s, 8.78, 0.70, 3.88, 5.62, fill="pale", line="pale")
    text(s, "RK3588", 9.06, 0.92, 2.25, 0.42, size=30, color="blue", bold=True, font="Aptos Display")
    text(s, "NPU\n实时检测\n系统", 9.05, 1.48, 2.75, 1.80, size=28, color="navy", bold=True)
    for i, label in enumerate(["YOLOv10", "RKNN", "RTSP", "RGA / INT8"]):
        pill(s, label, 9.08, 3.60 + i * 0.48, 1.36 if i < 3 else 1.70, color="cyan" if i != 1 else "orange")
    text(s, "基于嵌入式平台的\n目标检测系统研究", 0.78, 1.18, 6.8, 1.20, size=32, color="white", bold=True)
    text(s, "毕业设计答辩 | 5分钟报告 + 5分钟成果展示", 0.82, 2.54, 6.2, 0.30, size=15.5, color="pale2")
    text(s, "朱奕澄  ·  电子信息工程  ·  指导教师：滕国伟", 0.82, 5.80, 6.8, 0.25, size=13, color="pale2")
    text(s, "答辩时间：2026年5月18日 13:00", 0.82, 6.16, 5.0, 0.22, size=11.5, color="cyan")


def slide_goal(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "01 课题目标", "从“模型能检测”到“板端可演示系统”",
          "我的工作重点不是重新设计检测网络，而是把无人机检测模型稳定迁移到 RK3588 并完成实时视频链路。")
    node(s, "目标难点 1：小目标检测", "无人机在远距离画面中尺寸小，置信度容易波动，显示框也容易抖动。", 0.80, 1.75, 3.65, 1.05, "orange")
    node(s, "目标难点 2：端侧资源受限", "视频采集、预处理、NPU推理、后处理和推流会共同影响最终延迟。", 4.85, 1.75, 3.65, 1.05, "blue")
    node(s, "目标难点 3：演示必须稳定", "答辩现场优先实时画面，同时准备公开视频固定输入作为备用方案。", 8.90, 1.75, 3.65, 1.05, "green")
    text(s, "本文完成的系统目标", 0.86, 3.35, 3.0, 0.32, size=18, color="navy", bold=True)
    bullets(s, [
        "完成 YOLOv10 单类无人机模型到 ONNX / RKNN 的迁移与后处理适配",
        "实现 RK3588 板端 C++17 检测程序，支持固定视频、USB摄像头和 RTSP 输出",
        "围绕实时观看效果优化检测间隔、框平滑、动态 ROI、multi-context、RGA 与 INT8",
        "输出带框视频、检测 CSV、ROI JSONL 与软件告警记录，保证实验可复现"
    ], 0.95, 3.86, 11.5, 1.72, size=15.2)
    footer(s, 2)


def slide_architecture(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "02 系统总体方案", "视频输入 → 预处理 → NPU推理 → 后处理 → 输出展示",
          "系统按模块解耦，固定视频用于复现实验，实时 RTSP 用于答辩演示。")
    items = [
        ("视频输入", "MP4 / USB / RTSP", "cyan"),
        ("预处理", "resize / color / letterbox", "cyan"),
        ("NPU推理", "RKNN Runtime", "orange"),
        ("后处理", "置信度 + NMS", "blue"),
        ("策略控制", "N=3 / 平滑 / ROI", "green"),
        ("结果输出", "MP4 / RTSP / CSV", "red"),
    ]
    for i, (h, b, c) in enumerate(items):
        x = 0.62 + i * 2.08
        node(s, h, b, x, 2.02, 1.58, 0.86, c)
        if i < len(items) - 1:
            arrow(s, x + 1.58, 2.45, x + 1.94, 2.45)
    shape(s, 0.86, 4.10, 5.55, 1.42, fill="cream", line="line", rounded=True)
    text(s, "固定视频路径", 1.10, 4.30, 1.8, 0.24, size=15.5, color="orange", bold=True)
    bullets(s, ["输入一致、结果可复现", "输出带框视频 + CSV + ROI", "用于实验对比和备用演示"], 1.10, 4.68, 4.90, 0.58, size=12.3)
    shape(s, 6.90, 4.10, 5.55, 1.42, fill="pale", line="line", rounded=True)
    text(s, "实时演示路径", 7.14, 4.30, 1.8, 0.24, size=15.5, color="blue", bold=True)
    bullets(s, ["USB摄像头输入", "GStreamer RTSP推流", "电脑端同步观看告警画面"], 7.14, 4.68, 4.90, 0.58, size=12.3)
    footer(s, 3)


def slide_model(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "03 模型迁移与板端适配", "关键不是只转换格式，而是让输出张量与后处理一致",
          "最终稳定方案采用 FP RKNN；hybrid INT8 作为已验证但未默认启用的量化路径。")
    steps = [("best.pt", "训练权重"), ("best.onnx", "跨框架格式"), ("best.rknn", "板端模型"), ("C++ 后处理", "单类检测输出")]
    for i, (h, b) in enumerate(steps):
        x = 1.0 + i * 3.0
        shape(s, x, 1.85, 1.95, 0.88, fill="pale2", line="blue", rounded=True)
        text(s, h, x + 0.08, 2.05, 1.79, 0.16, size=14.5, color="navy", bold=True, align=PP_ALIGN.CENTER)
        text(s, b, x + 0.08, 2.32, 1.79, 0.16, size=10.8, color="muted", align=PP_ALIGN.CENTER)
        if i < 3:
            arrow(s, x + 1.95, 2.29, x + 2.70, 2.29)
    table(s, ["问题", "处理方式", "答辩要点"], [
        ("end2end 输出不兼容", "改用 end2end=False", "输出保持 1×5×8400，C++端解析"),
        ("单类无人机检测", "只保留 drone 类", "后处理重点在置信度阈值和坐标还原"),
        ("full INT8 置信度塌缩", "尝试 hybrid INT8", "恢复检测能力，但默认仍选择 FP 稳定方案"),
    ], 0.95, 3.58, [2.45, 3.45, 5.55], row_h=0.54, font_size=10)
    footer(s, 4)


def slide_realtime(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "04 实时链路优化", "最终演示配置优先保证低延迟、连续画面和稳定检测框",
          "实时画面不是单纯追求每帧检测，而是要让老师看到“画面不卡、框不乱跳、告警清楚”。")
    node(s, "检测间隔 N=3", "降低 NPU 调用频率，减少实时链路积压。", 0.80, 1.78, 3.45, 1.05, "blue")
    node(s, "框平滑 box_smooth", "抑制相邻检测框突变，提高观看稳定性。", 4.95, 1.78, 3.45, 1.05, "green")
    node(s, "运动预测 / 光流", "在非完整检测帧中估计位置，提升连续性。", 9.10, 1.78, 3.45, 1.05, "orange")
    text(s, "多线程流水线", 0.86, 3.32, 2.0, 0.28, size=18, color="navy", bold=True)
    table(s, ["线程/模块", "作用", "为什么重要"], [
        ("采集线程", "持续读取最新帧", "避免推理阻塞摄像头输入"),
        ("推理线程", "调用 RKNN/NPU", "隔离模型耗时，便于扩展 context"),
        ("输出线程", "绘制检测框并推流", "保证电脑端能同步观看"),
    ], 0.95, 3.82, [2.2, 4.4, 4.8], row_h=0.52, font_size=10.4)
    footer(s, 5)


def slide_hardware(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "05 硬件优化探索", "NPU 多 context、RGA 和 INT8 都做了验证，但默认方案选择稳定性",
          "答辩时要强调：硬件优化不是没有做，而是根据实验结果做了工程取舍。")
    table(s, ["方向", "完成情况", "结论"], [
        ("multi-context NPU", "已实现多 RKNN context 并行", "逐帧检测吞吐提高，但 context 过多会增加调度压力"),
        ("RGA 预处理", "已实现可切换 RGA resize/color/letterbox", "功能验证通过，默认仍保留 OpenCV 稳定方案"),
        ("INT8 量化", "full INT8 与 hybrid INT8 均验证", "full INT8 置信度塌缩；hybrid INT8 恢复检测但未稳定超过 FP"),
        ("zero-copy 输入", "已完成阶段性验证", "真正收益依赖 MPP/RGA/NPU 的完整零拷贝链路"),
    ], 0.80, 1.70, [2.2, 4.4, 5.55], row_h=0.68, font_size=10.2)
    text(s, "答辩回答口径", 0.96, 5.30, 2.2, 0.25, size=16, color="navy", bold=True)
    bullets(s, [
        "RK3588 的瓶颈不只在 NPU 计算，还包括输入设置、内存搬运、编码推流和队列调度。",
        "因此最终演示选择 FP 稳定方案，同时保留 RGA / INT8 / zero-copy 的实验记录作为优化探索。"
    ], 1.02, 5.72, 11.5, 0.78, size=13.2)
    footer(s, 6)


def slide_results(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "06 实验结果", "公开无人机视频和板端测试验证了系统可用性",
          "这里建议讲 40 秒：先讲检测结果，再讲稳定性，再落到推荐配置。")
    add_image(s, IMG_DEMO, 0.75, 1.55, 6.10, 3.43)
    text(s, "板端公开视频固定输入结果", 0.75, 5.10, 6.10, 0.25, size=12.2, color="muted", align=PP_ALIGN.CENTER)
    metric(s, "1200", "测试帧数", 7.30, 1.72, color="blue")
    metric(s, "397", "有检测帧", 9.25, 1.72, color="green")
    metric(s, "404", "检测框数量", 11.15, 1.72, color="orange")
    metric(s, "81.27ms", "平均推理耗时", 7.30, 3.12, color="blue")
    metric(s, "24,700", "1小时 NPU 推理次数", 9.25, 3.12, color="green")
    metric(s, "0", "崩溃 / 明显泄漏", 11.15, 3.12, color="red")
    bullets(s, [
        "公开视频固定输入：输出带框视频、CSV、ROI JSONL 和告警日志。",
        "长时间压力测试：1小时独立推理无崩溃，RSS 内存保持稳定。",
        "最终推荐：实时演示使用 FP RKNN + 检测间隔 + 框平滑；公开视频作为备用验证。"
    ], 7.25, 4.80, 5.30, 1.10, size=12.5)
    footer(s, 7)


def slide_demo(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "07 成果展示", "两条验证链路：实时检测与固定视频复现",
          "现场展示关注端到端可运行，备用链路保证结果可重复、可对比。")

    text(s, "实时检测链路", 0.86, 1.78, 2.2, 0.24, size=16.5, color="blue", bold=True)
    text(s, "证明系统可以接入真实摄像头，并把检测结果低延迟推到电脑端。", 2.08, 1.81, 8.5, 0.22, size=10.8, color="muted")
    node(s, "USB 摄像头", "现场画面输入", 0.90, 2.20, 2.05, 0.82, "blue", fill="pale")
    node(s, "RK3588 / NPU", "实时检测 + 轻量跟踪", 3.45, 2.20, 2.35, 0.82, "green", fill="pale")
    node(s, "RTSP 推流", "GStreamer 输出", 6.30, 2.20, 2.05, 0.82, "blue", fill="pale")
    node(s, "电脑端观看", "检测框 + UAV ALERT", 8.85, 2.20, 2.45, 0.82, "orange", fill="pale")
    arrow(s, 2.98, 2.61, 3.38, 2.61, "blue")
    arrow(s, 5.83, 2.61, 6.23, 2.61, "blue")
    arrow(s, 8.38, 2.61, 8.78, 2.61, "blue")

    shape(s, 0.90, 3.38, 10.40, 0.02, fill="line", line="line")

    text(s, "固定视频复现链路", 0.86, 3.82, 2.65, 0.24, size=16.5, color="orange", bold=True)
    text(s, "证明同一模型和后处理逻辑可以在固定输入下反复验证。", 2.42, 3.85, 8.5, 0.22, size=10.8, color="muted")
    node(s, "公开视频", "统一输入源", 0.90, 4.24, 2.05, 0.82, "orange", fill="cream")
    node(s, "rk_yolo_video", "板端离线推理", 3.45, 4.24, 2.35, 0.82, "blue", fill="cream")
    node(s, "检测产物", "带框视频 / CSV / ROI", 6.30, 4.24, 2.35, 0.82, "green", fill="cream")
    node(s, "结果复核", "可记录、可对比", 9.15, 4.24, 2.15, 0.82, "orange", fill="cream")
    arrow(s, 2.98, 4.65, 3.38, 4.65, "orange")
    arrow(s, 5.83, 4.65, 6.23, 4.65, "orange")
    arrow(s, 8.68, 4.65, 9.08, 4.65, "orange")

    shape(s, 0.92, 5.72, 10.35, 0.58, fill="navy", line="navy", rounded=True)
    text(s, "展示策略：优先展示实时画面；若现场光照、网络或摄像头状态不稳定，立即切换到固定视频结果，仍能完整说明模型、NPU 推理、后处理和告警输出。", 1.15, 5.90, 9.86, 0.23, size=12.4, color="white")
    footer(s, 8)


def slide_summary(prs, blank):
    s = prs.slides.add_slide(blank)
    title(s, "08 工作总结", "我的主要贡献：把算法变成可运行、可展示、可复现的板端系统",
          "最后 30 秒只讲亮点，不再展开技术细节。")
    node(s, "1. 模型迁移", "完成 PyTorch / ONNX / RKNN 转换，解决输出张量与后处理适配问题。", 0.95, 1.72, 3.72, 1.18, "blue")
    node(s, "2. 实时系统", "实现固定视频与 RTSP 实时链路，支持带框画面、告警和日志输出。", 4.82, 1.72, 3.72, 1.18, "green")
    node(s, "3. 性能优化", "验证 multi-context、RGA、INT8、zero-copy，并形成稳定演示配置。", 8.68, 1.72, 3.72, 1.18, "orange")
    shape(s, 1.35, 3.55, 10.65, 1.62, fill="navy", line="navy", rounded=True)
    text(s, "最终结论", 1.70, 3.88, 1.5, 0.28, size=17, color="cyan", bold=True)
    text(s, "RK3588 能支撑无人机检测模型的端侧部署，但最终演示配置不能只依据 NPU 单次推理速度选择，必须同时考虑输入搬运、推流延迟、检测间隔、框稳定性和长时间运行可靠性。", 3.15, 3.83, 8.45, 0.75, size=15.3, color="white")
    text(s, "谢谢各位老师，请批评指正。", 0.90, 6.03, 11.6, 0.42, size=24, color="navy", bold=True, align=PP_ALIGN.CENTER)
    footer(s, 9)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide_cover(prs, blank)
    slide_goal(prs, blank)
    slide_architecture(prs, blank)
    slide_model(prs, blank)
    slide_realtime(prs, blank)
    slide_hardware(prs, blank)
    slide_results(prs, blank)
    slide_demo(prs, blank)
    slide_summary(prs, blank)
    prs.save(OUT_PPTX)

    notes = """# 朱奕澄+答辩PPT 5分钟讲稿

## 第1页 封面（约20秒）
各位老师同学好，我是朱奕澄。我的毕业设计题目是《基于嵌入式平台的目标检测系统研究》。本课题主要围绕 RK3588 开发板，把无人机检测模型部署到板端，并完成实时视频显示和性能优化。

## 第2页 课题目标（约40秒）
这个课题不是单纯验证模型能不能检测无人机，而是让它在 RK3588 上形成一套可运行、可展示、可记录的系统。实际问题包括小目标检测不稳定、板端算力和内存搬运受限、视频推流容易产生延迟等。因此我的工作重点是模型迁移、C++板端程序、多线程实时链路和硬件优化验证。

## 第3页 系统总体方案（约40秒）
系统分为视频输入、预处理、NPU推理、后处理、策略控制和输出展示。固定视频路径用于可重复实验，实时 RTSP 路径用于现场演示。输出不仅有带框视频，还有 CSV、ROI JSONL 和告警日志，便于后续分析。

## 第4页 模型迁移（约40秒）
模型从 best.pt 导出为 ONNX，再转换为 RKNN。早期遇到的关键问题是输出张量与后处理不匹配，后来采用 end2end=False，使输出保持为 1×5×8400，并在 C++ 端完成置信度筛选、坐标还原和 NMS。最终演示采用 FP RKNN 稳定方案。

## 第5页 实时链路优化（约45秒）
实时演示不能只追求每帧检测。为了降低延迟和抖动，系统加入检测间隔、框平滑、运动预测和动态 ROI。采集、推理、输出线程解耦后，摄像头输入不会被 NPU 推理直接阻塞，电脑端可以通过 RTSP 同步观看。

## 第6页 硬件优化探索（约45秒）
我验证了 multi-context、RGA、INT8 和 zero-copy。multi-context 能提升逐帧推理吞吐，但 context 过多会带来调度和延迟问题。RGA 和 zero-copy 已完成阶段性验证，但当前稳定演示仍选择 OpenCV 预处理。full INT8 出现置信度塌缩，hybrid INT8 能恢复检测，但未稳定超过 FP 方案。

## 第7页 实验结果（约50秒）
公开视频固定输入测试中，板端处理 1200 帧，其中 397 帧有检测，共产生 404 个检测框，平均推理耗时约 81.27 ms。长时间测试中，系统完成 24700 次 NPU 推理，无崩溃和明显内存泄漏。这说明系统具备稳定运行和演示能力。

## 第8页 成果展示（约35秒）
这一页不用逐条读，主要说明我准备了两条验证链路。现场优先展示实时摄像头：USB 摄像头接入 RK3588，检测结果经 RTSP 推到电脑端观看，并显示检测框和 UAV ALERT 告警条。如果现场光照、网络或摄像头状态不理想，就切到固定视频结果。固定视频和实时画面使用同一套模型、后处理和告警逻辑，所以即使切换备用方案，也能完整展示系统能力。

## 第9页 总结（约30秒）
总结来说，我完成了模型迁移、板端实时检测系统和多种硬件优化验证。最终结论是，RK3588 能够支持无人机检测模型的端侧部署，但实时演示配置需要综合考虑检测稳定性、推流延迟、输入搬运和长时间运行可靠性。我的汇报结束，谢谢各位老师。
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")
    print(f"PPTX={OUT_PPTX}")
    print(f"NOTES={OUT_NOTES}")
    print(f"SLIDES={len(prs.slides)}")


if __name__ == "__main__":
    build()
