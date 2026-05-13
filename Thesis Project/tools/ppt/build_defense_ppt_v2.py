from __future__ import annotations

from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[2]
VENDOR = ROOT.parent.parent / ".vendor"
if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "paper" / "defense_ppt"
OUT_PPTX = OUT_DIR / "defense_presentation_v2.pptx"
OUT_NOTES = OUT_DIR / "defense_presentation_speaker_notes_v2.md"

W, H = Inches(13.333), Inches(7.5)

COLORS = {
    "navy": RGBColor(13, 31, 48),
    "deep": RGBColor(20, 45, 70),
    "blue": RGBColor(24, 107, 176),
    "cyan": RGBColor(0, 150, 196),
    "teal": RGBColor(28, 150, 136),
    "green": RGBColor(69, 165, 95),
    "orange": RGBColor(230, 130, 38),
    "red": RGBColor(205, 62, 62),
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(99, 116, 139),
    "pale": RGBColor(241, 247, 252),
    "panel": RGBColor(225, 238, 248),
    "line": RGBColor(190, 207, 222),
    "white": RGBColor(255, 255, 255),
}


def set_font(run, size=18, color="ink", bold=False, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = COLORS[color]
    run.font.bold = bold


def add_text(slide, value, x, y, w, h, size=18, color="ink", bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    set_font(run, size=size, color=color, bold=bold, font=font)
    return box


def add_lines(slide, lines, x, y, w, h, size=15.5, color="ink", bullet=True, gap=3):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        set_font(run, size=size, color=color)
    return box


def rect(slide, x, y, w, h, fill="pale", line="line", radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = COLORS[fill]
    s.line.color.rgb = COLORS[line]
    s.line.width = Pt(0.8)
    return s


def chip(slide, value, x, y, w, fill="blue", size=9.5):
    rect(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    add_text(slide, value, x + 0.05, y + 0.065, w - 0.1, 0.14, size, "white", True, PP_ALIGN.CENTER)


def node(slide, title, body, x, y, w, h, accent="blue", fill="white"):
    rect(slide, x, y, w, h, fill=fill, line="line", radius=True)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.28, 0.22, 13.5, "navy", True)
    add_text(slide, body, x + 0.18, y + 0.48, w - 0.28, h - 0.52, 10.5, "muted")


def arrow(slide, x1, y1, x2, y2, color="blue", width=1.8):
    a = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    a.line.color.rgb = COLORS[color]
    a.line.width = Pt(width)
    return a


def add_title(slide, section, title, subtitle=None):
    add_text(slide, section, 0.66, 0.28, 3.2, 0.25, 10.5, "cyan", True)
    add_text(slide, title, 0.64, 0.58, 9.8, 0.52, 27, "navy", True)
    if subtitle:
        add_text(slide, subtitle, 0.66, 1.12, 10.7, 0.26, 12.5, "muted")


def add_footer(slide, page):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.64), Inches(6.9), Inches(12.05), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()
    add_text(slide, "基于嵌入式平台的目标检测系统研究", 0.64, 7.02, 5.0, 0.18, 8.2, "muted")
    add_text(slide, f"{page:02d}", 12.18, 7.02, 0.45, 0.18, 8.5, "muted", align=PP_ALIGN.RIGHT)


def add_table(slide, headers, rows, x, y, col_w, row_h=0.52, header_fill="blue"):
    for i, head in enumerate(headers):
        rect(slide, x + sum(col_w[:i]), y, col_w[i], 0.42, fill=header_fill, line=header_fill)
        add_text(slide, head, x + sum(col_w[:i]) + 0.06, y + 0.1, col_w[i] - 0.12, 0.14, 9.8, "white", True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        fill = "pale" if r % 2 == 0 else "white"
        yy = y + 0.44 + r * row_h
        for i, cell in enumerate(row):
            xx = x + sum(col_w[:i])
            rect(slide, xx, yy, col_w[i], row_h - 0.02, fill=fill, line="line")
            add_text(slide, cell, xx + 0.07, yy + 0.12, col_w[i] - 0.14, 0.18, 9.8, "ink", bold=(cell in ["已完成", "已验证", "实验路径"]))


def add_cover(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, fill="navy", line="navy")
    rect(s, 8.1, 0, 5.25, 7.5, fill="blue", line="blue")
    rect(s, 8.78, 0.62, 3.72, 5.6, fill="pale", line="pale")
    add_text(s, "RK3588", 9.08, 0.94, 2.3, 0.46, 30, "blue", True, font="Aptos Display")
    add_text(s, "NPU 多线程\n实时目标检测", 9.08, 1.55, 2.9, 1.0, 25, "navy", True)
    for i, label in enumerate(["PyTorch", "ONNX", "RKNN", "RTSP"]):
        chip(s, label, 9.08, 3.02 + i * 0.50, 1.15, fill="teal" if label != "RKNN" else "orange")
    add_text(s, "基于嵌入式平台的\n目标检测系统研究", 0.76, 1.06, 6.85, 1.05, 29, "white", True)
    add_text(s, "本科毕业论文（设计）答辩汇报", 0.80, 2.12, 4.8, 0.28, 16.5, "panel")
    add_text(s, "朱奕澄  |  通信与信息工程学院  |  电子信息工程", 0.82, 5.85, 6.8, 0.24, 12.5, "panel")
    add_text(s, "研究重点：模型迁移、RKNN部署、实时流水线、NPU多 context 并行、RGA 与 INT8 工程验证", 0.82, 6.26, 7.4, 0.25, 11.5, "teal")


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    add_cover(prs, blank)

    s = prs.slides.add_slide(blank)
    add_title(s, "01 研究背景", "边缘端小目标检测需要系统级优化", "无人机目标往往小、远、受背景干扰；在 RK3588 上还要面对实时性和硬件资源约束。")
    add_lines(s, [
        "远距离无人机在图像中占比很小，容易受天空、建筑、树木、鸟类和光照变化影响。",
        "嵌入式端的瓶颈不只在模型推理，还包括图像预处理、内存拷贝、后处理、显示和视频编码。",
        "本课题重点不是提出全新检测网络，而是完成从训练模型到 RK3588 板端实时运行的工程闭环。"
    ], 0.85, 1.72, 6.0, 1.9)
    node(s, "检测难点", "小目标、背景干扰、框抖动、漏检与误检", 7.45, 1.55, 4.8, 0.85, "blue")
    node(s, "部署难点", "NPU 推理、CPU 后处理、内存带宽、RTSP 输出", 7.45, 2.65, 4.8, 0.85, "teal")
    node(s, "研究定位", "面向 RK3588 的目标检测系统移植、优化与验证", 7.45, 3.75, 4.8, 0.85, "orange")
    add_footer(s, 2)

    s = prs.slides.add_slide(blank)
    add_title(s, "02 任务对照", "任务书要求与当前完成情况", "主动把“已完成、实验路径、替代实现”的边界讲清楚，会比回避问题更专业。")
    add_table(s, ["任务线", "状态", "答辩口径"], [
        ("模型迁移", "已完成", "PyTorch/ONNX/RKNN 链路跑通，FP RKNN 作为稳定主线"),
        ("实时检测", "已完成", "固定视频与 USB 摄像头 RTSP 均可输出带框画面"),
        ("NPU 多线程", "已验证", "multi-context 用于证明 NPU 并行推理能力"),
        ("RGA 预处理", "实验路径", "strict RGA 可强制启用并板端验证，不虚称默认最优"),
        ("INT8 量化", "实验路径", "full INT8 发现小目标置信度问题，hybrid INT8 恢复检测"),
        ("闭环报警", "替代实现", "软件报警叠加、alarm CSV 与 GPIO 兼容接口")
    ], 0.78, 1.55, [1.75, 1.55, 8.65], row_h=0.58)
    add_footer(s, 3)

    s = prs.slides.add_slide(blank)
    add_title(s, "03 总体架构", "视频输入到检测输出的分层流水线", "同一套检测器支撑固定视频实验和摄像头实时演示。")
    items = [
        ("视频输入", "MP4 / USB / RTSP", "teal"),
        ("预处理", "OpenCV / RGA 可切换", "cyan"),
        ("NPU 推理", "RKNN Runtime / context", "orange"),
        ("后处理", "解码 / 阈值 / NMS", "blue"),
        ("策略控制", "间隔检测 / ROI / 跟踪", "green"),
        ("结果输出", "MP4 / CSV / RTSP / 报警", "red"),
    ]
    for i, (t, b, c) in enumerate(items):
        x = 0.65 + i * 2.05
        node(s, t, b, x, 2.0, 1.52, 0.86, c)
        if i < len(items) - 1:
            arrow(s, x + 1.52, 2.43, x + 1.92, 2.43)
    add_lines(s, [
        "固定视频路径：负责可重复实验，输出带框视频、检测 CSV、ROI JSONL 和 alarm_events.csv。",
        "实时 RTSP 路径：负责 USB 摄像头演示，电脑端可直接观看实时检测画面。",
        "硬件优化路径：multi-context、RGA、zero-copy、profiling 均通过环境变量开关对比。"
    ], 1.05, 4.15, 11.2, 1.25, 14.5)
    add_footer(s, 4)

    s = prs.slides.add_slide(blank)
    add_title(s, "04 模型迁移", "从训练权重到 RK3588 可执行模型", "迁移重点是输出张量与后处理一致，而不是只完成格式转换。")
    steps = [("best.pt", "训练权重"), ("best.onnx", "中间格式"), ("best.rknn", "板端模型"), ("C++ decode", "统一后处理")]
    for i, (t, b) in enumerate(steps):
        x = 1.0 + i * 3.0
        node(s, t, b, x, 1.8, 2.05, 0.86, ["blue", "cyan", "orange", "green"][i])
        if i < 3:
            arrow(s, x + 2.05, 2.23, x + 2.65, 2.23)
    rect(s, 1.08, 3.55, 11.12, 1.25, fill="pale", line="line", radius=True)
    add_text(s, "输出张量说明", 1.34, 3.80, 1.6, 0.22, 14, "navy", True)
    add_text(s, "8400 = 80×80 + 40×40 + 20×20，来自 YOLOv10 三个检测头；当前单类模型为 [1,5,8400]，后处理同时兼容 COCO 84 通道和 end2end 6 通道。", 2.95, 3.76, 8.85, 0.38, 13.5, "ink")
    add_lines(s, ["稳定演示主线采用 FP RKNN；INT8 作为量化优化实验线，保留完整问题分析。"], 1.2, 5.25, 10.2, 0.42, 14.5)
    add_footer(s, 5)

    s = prs.slides.add_slide(blank)
    add_title(s, "05 RK3588 平台", "算力不是单点资源，而是异构协同", "CPU、NPU、RGA、视频编解码与内存带宽共同决定端到端效果。")
    metrics = [("6 TOPS", "NPU INT8 算力"), ("4×A76 + 4×A55", "big.LITTLE CPU"), ("Mali-G610", "GPU 图形能力"), ("LPDDR4/5", "内存与带宽约束")]
    for i, (v, l) in enumerate(metrics):
        x = 0.88 + i * 3.05
        rect(s, x, 1.72, 2.45, 1.15, fill="pale", line="line", radius=True)
        add_text(s, v, x + 0.18, 1.98, 2.05, 0.32, 19, ["blue", "teal", "orange", "green"][i], True)
        add_text(s, l, x + 0.18, 2.42, 2.05, 0.16, 10.5, "muted")
    add_lines(s, [
        "NPU 负责神经网络前向推理，但视频解码、图像预处理、后处理和推流仍会占用 CPU 与内存带宽。",
        "RGA 用于颜色转换、缩放和 letterbox 等图像处理，目标是减少 CPU 预处理压力。",
        "因此论文实验采用 profiling 拆分 prepare、input、run、output、decode、render 等阶段。"
    ], 1.05, 3.75, 10.9, 1.35, 14.5)
    add_footer(s, 6)

    s = prs.slides.add_slide(blank)
    add_title(s, "06 实时流水线与 NPU 多 context", "把阻塞串行过程拆成可并行模块", "导师关注的“多线程方向”主要落在任务队列、独立 RKNN context 和结果发布策略上。")
    node(s, "采集线程", "读取最新帧，写入有限队列", 0.85, 2.0, 1.85, 0.85, "teal")
    node(s, "任务队列", "FrameTask / BoundedQueue", 3.05, 2.0, 1.95, 0.85, "blue")
    node(s, "Worker 1", "RKNN Context 1", 5.35, 1.45, 1.9, 0.72, "orange")
    node(s, "Worker 2", "RKNN Context 2", 5.35, 2.35, 1.9, 0.72, "orange")
    node(s, "结果队列", "DetectionResult", 7.95, 2.0, 1.9, 0.85, "green")
    node(s, "显示/推流", "画框、RTSP、日志", 10.4, 2.0, 1.95, 0.85, "red")
    arrow(s, 2.7, 2.43, 3.03, 2.43)
    arrow(s, 5.0, 2.43, 5.32, 1.8)
    arrow(s, 5.0, 2.43, 5.32, 2.72)
    arrow(s, 7.25, 1.8, 7.93, 2.43)
    arrow(s, 7.25, 2.72, 7.93, 2.43)
    arrow(s, 9.85, 2.43, 10.38, 2.43)
    add_lines(s, [
        "每个 worker 持有独立 RKNN context，避免多个线程共享同一推理上下文造成同步冲突。",
        "多 context 在每帧检测场景中更能体现 NPU 并行能力；实时观看则结合检测间隔和跟踪策略。"
    ], 1.05, 4.72, 10.8, 0.86, 14.2)
    add_footer(s, 7)

    s = prs.slides.add_slide(blank)
    add_title(s, "07 RGA 硬件预处理", "已实现、可强制启用、可与 OpenCV 对比", "RGA 部分按任务书要求形成“视频采集 -> RGA 预处理 -> NPU 推理 -> 后处理”的实验路径。")
    add_table(s, ["路径", "实现内容", "验证口径"], [
        ("rga_resize", "RGA resize", "可运行对比"),
        ("rga_cvt_resize", "BGR->RGB + resize", "减少 CPU 颜色转换"),
        ("rga_letterbox", "缩放 + padding", "更接近模型输入预处理"),
        ("strict RGA", "RK_YOLO_REQUIRE_RGA=1", "失败即退出，不回退 OpenCV")
    ], 0.86, 1.65, [2.0, 4.2, 5.2], row_h=0.64, header_fill="teal")
    add_lines(s, [
        "板端日志可确认 rga_api 版本和 require_rga=on，说明不是仅停留在论文描述。",
        "RGA 是否作为默认主路径取决于端到端延迟，不以单个阶段变快作为最终结论。"
    ], 1.0, 5.0, 10.6, 0.72, 14.2)
    add_footer(s, 8)

    s = prs.slides.add_slide(blank)
    add_title(s, "08 INT8 量化实验", "full INT8 暴露问题，hybrid INT8 恢复检测", "这一页建议答辩时主动讲，避免老师追问时显得被动。")
    add_table(s, ["方案", "现象", "结论"], [
        ("FP RKNN", "检测稳定，作为最终演示主线", "稳定优先"),
        ("full INT8", "可运行但小目标置信度不稳定，出现不出框", "发现量化敏感点"),
        ("hybrid INT8", "保护关键输出层精度后恢复检测", "更接近工程可用"),
        ("后续方向", "扩大校准集、逐层误差分析、精度/速度闭环", "继续优化")
    ], 0.86, 1.62, [2.05, 5.05, 4.35], row_h=0.66, header_fill="orange")
    rect(s, 1.05, 5.08, 10.9, 0.62, fill="panel", line="line", radius=True)
    add_text(s, "答辩口径：本文完成了 INT8 链路、板端运行验证和量化误差分析；最终演示选择 FP RKNN，是为了保证系统稳定性。", 1.28, 5.27, 10.35, 0.18, 12.8, "navy", True)
    add_footer(s, 9)

    s = prs.slides.add_slide(blank)
    add_title(s, "09 实时检测策略", "检测间隔、ROI、轻量跟踪与框平滑", "目标是在不明显牺牲观看体验的情况下减少 NPU 调用和画面抖动。")
    cards = [
        ("detect_every_n", "每 N 帧执行一次完整检测，中间帧复用或跟踪结果。"),
        ("动态 ROI", "围绕历史目标区域裁剪，并周期性全帧刷新。"),
        ("轻量跟踪", "Motion 线性预测；OpticalFlow 使用 Lucas-Kanade 稀疏光流。"),
        ("框平滑", "对检测框坐标进行加权融合，减小左右抖动。")
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.85 + (i % 2) * 6.0
        y = 1.68 + (i // 2) * 1.55
        node(s, t, b, x, y, 4.95, 0.98, ["blue", "teal", "orange", "green"][i])
    rect(s, 1.08, 5.25, 10.75, 0.54, fill="pale", line="line", radius=True)
    add_text(s, "推荐演示配置：detect_every_n=2，motion tracking，box smoothing on，dynamic ROI on。", 1.32, 5.42, 10.2, 0.16, 13, "navy", True)
    add_footer(s, 10)

    s = prs.slides.add_slide(blank)
    add_title(s, "10 检测效果与报警可视化", "离线视频和实时 RTSP 都支持画面级报警提示", "没有外接继电器/蜂鸣器时，软件报警更适合现场演示和论文记录。")
    img = ROOT / "eval_runs" / "alarm_validation" / "fig1_alarm_overlay_frame1.jpg"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0.95), Inches(1.65), width=Inches(5.9))
    else:
        rect(s, 0.95, 1.65, 5.9, 3.3, fill="deep", line="deep")
        add_text(s, "检测画面截图", 2.95, 3.1, 1.5, 0.2, 14, "white", True)
    add_lines(s, [
        "离线视频：输出带框 MP4、逐帧 CSV、ROI JSONL 与 alarm_events.csv。",
        "实时 RTSP：画面顶部显示 UAV ALERT / NORMAL 状态，日志打印 alarm=on/off。",
        "该替代方案可直观看到“识别到目标 -> 触发报警”的闭环过程。"
    ], 7.25, 1.9, 4.85, 1.65, 14.2)
    add_footer(s, 11)

    s = prs.slides.add_slide(blank)
    add_title(s, "11 实验结果与配置推荐", "系统级最优与任务级最优要分开讲", "多 context 证明硬件并行能力；实时演示更关注端到端流畅和稳定。")
    rows = [
        ("Baseline", "稳定基线", "用于对照"),
        ("Multi-context", "NPU FPS 提升", "硬件并行证据"),
        ("Policy N=2/N=3", "降低调用频率", "实时观看默认方向"),
        ("RGA/zero-copy", "部分阶段有收益", "实验路径")
    ]
    for i, (a, b, c) in enumerate(rows):
        x = 0.86 + i * 3.05
        rect(s, x, 1.88, 2.48, 2.25, fill="pale" if i != 1 else "panel", line="line", radius=True)
        chip(s, a, x + 0.22, 2.15, 1.65, fill=["blue", "orange", "green", "teal"][i])
        add_text(s, b, x + 0.22, 2.78, 2.05, 0.25, 15.5, "navy", True)
        add_text(s, c, x + 0.22, 3.30, 2.05, 0.28, 12.0, "muted")
    add_lines(s, [
        "系统级推荐：FP RKNN + 策略优化，用于稳定演示和实时观看。",
        "任务级推荐：保留 RGA、INT8、multi-context 对比实验，体现硬件优化研究过程。"
    ], 1.05, 5.05, 10.7, 0.72, 14.5)
    add_footer(s, 12)

    s = prs.slides.add_slide(blank)
    add_title(s, "12 工作总结与不足展望", "完成工程闭环，也保留真实边界", "答辩时最重要的是把贡献和限制都讲清楚。")
    node(s, "系统闭环", "训练/转换/部署/固定视频/RTSP 实时显示", 0.95, 1.75, 3.5, 1.05, "teal", "pale")
    node(s, "优化闭环", "检测间隔、跟踪、ROI、multi-context、RGA、INT8", 4.95, 1.75, 3.5, 1.05, "orange", "pale")
    node(s, "论文闭环", "实验记录、代码映射、格式合规、附录材料", 8.95, 1.75, 3.5, 1.05, "green", "pale")
    add_lines(s, [
        "不足：full INT8 仍需更系统的校准集和精度闭环；RGA 与 zero-copy 需继续优化端到端收益。",
        "不足：暂未接入真实继电器/蜂鸣器/云台，真实飞行场景测试受条件限制。",
        "展望：接入物理报警或云台接口，扩大真实无人机场景，完善 INT8 与 RGA 主路径。"
    ], 1.05, 4.08, 10.9, 1.2, 14.2)
    add_footer(s, 13)

    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, fill="navy", line="navy")
    add_text(s, "谢谢各位老师", 0.88, 1.08, 4.5, 0.52, 34, "white", True)
    add_text(s, "欢迎批评指正", 0.9, 1.86, 3.0, 0.32, 18, "panel")
    rect(s, 0.9, 3.12, 10.9, 1.0, fill="deep", line="deep", radius=True)
    add_text(s, "答辩主线：本课题完成了一个可运行、可测量、可解释的 RK3588 嵌入式目标检测系统。", 1.18, 3.48, 10.2, 0.25, 20, "teal", True)
    add_text(s, "关键词：RK3588  |  YOLOv10  |  RKNN  |  NPU 多线程  |  RGA  |  INT8", 1.0, 5.85, 8.6, 0.22, 12.5, "panel")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    return len(prs.slides)


NOTES = """# 答辩 PPT 讲稿提示（v2）

## 1. 封面
一句话开场：本课题不是单纯训练一个模型，而是完成目标检测算法从 PC 训练到 RK3588 板端实时运行的移植、优化和验证。

## 2. 研究背景
强调小目标和边缘端实时性的双重困难，自然引出 NPU 多线程和硬件优化。

## 3. 任务对照
主动说明“已完成、实验路径、替代实现”的边界。RGA 和 INT8 是工程验证线，FP RKNN 是稳定演示线。

## 4. 系统架构
按视频输入、预处理、RKNN 推理、后处理、策略控制、输出六层讲。

## 5. 模型迁移
讲清楚 PyTorch -> ONNX -> RKNN，以及为什么后处理必须兼容不同输出张量。

## 6. RK3588 平台
强调 NPU、CPU、RGA、视频编解码和内存带宽共同影响实时性。

## 7. 多线程与多 context
这是答辩重点。每个 worker 持有独立 RKNN context，避免共享上下文同步冲突。

## 8. RGA
口径：已实现 strict RGA 路径并板端验证，但是否默认启用取决于端到端收益。

## 9. INT8
口径：full INT8 暴露小目标置信度问题，hybrid INT8 恢复检测；最终演示 FP RKNN 保证稳定性。

## 10. 实时策略
解释框抖动如何通过检测间隔、运动预测、光流候选、动态 ROI 和平滑缓解。

## 11. 检测效果
现场可切换到实时 RTSP 或检测视频演示。强调软件报警是无外设条件下的可视闭环。

## 12. 实验推荐
multi-context 是硬件并行证据；策略优化是实时观看默认方案；RGA/INT8 是硬件优化研究线。

## 13. 总结
用系统闭环、优化闭环、论文闭环收束。最后强调可运行、可测量、可解释。
"""


if __name__ == "__main__":
    count = build_deck()
    OUT_NOTES.write_text(NOTES, encoding="utf-8")
    print(f"created {OUT_PPTX} ({count} slides)")
    print(f"created {OUT_NOTES}")
